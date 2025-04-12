from pathlib import Path
from PyPDF2 import PdfReader, PdfWriter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from io import BytesIO

# === Thiết lập đường dẫn ===
input_dir = Path("raw/watermarked")
output_dir = Path("raw/files")
output_dir.mkdir(parents=True, exist_ok=True)

# === Watermark PDF ===
def create_pdf_watermark(text):
    packet = BytesIO()
    c = canvas.Canvas(packet, pagesize=letter)
    c.setFont("Helvetica", 30)
    c.setFillColorRGB(0.6, 0.6, 0.6, alpha=0.3)
    c.rotate(45)
    c.drawString(100, 0, text)
    c.save()
    packet.seek(0)
    return PdfReader(packet)

def watermark_pdf(input_path, output_path, watermark_text):
    watermark = create_pdf_watermark(watermark_text)
    writer = PdfWriter()
    reader = PdfReader(input_path)
    for page in reader.pages:
        page.merge_page(watermark.pages[0])
        writer.add_page(page)
    with open(output_path, "wb") as f:
        writer.write(f)

# === Watermark PPTX ===
def watermark_pptx(input_path, output_path, watermark_text):
    prs = Presentation(input_path)
    for slide in prs.slides:
        left = Inches(1)
        top = Inches(5)
        width = Inches(8)
        height = Inches(1)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = watermark_text
        p = text_frame.paragraphs[0]
        run = p.runs[0]
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = RGBColor(200, 200, 200)
    prs.save(output_path)

# === Xử lý file ===
watermark_text = "NhanNgocThach – For Educational Use Only"
processed_files = []

for file in input_dir.iterdir():
    out = output_dir / file.name

    # Bỏ qua nếu đã có
    if out.exists():
        print(f"⏩ Bỏ qua: {file.name} (đã có watermark)")
        continue

    try:
        if file.suffix.lower() == ".pdf":
            watermark_pdf(file, out, watermark_text)
            processed_files.append(out.name)
        elif file.suffix.lower() == ".pptx":
            watermark_pptx(file, out, watermark_text)
            processed_files.append(out.name)
    except Exception as e:
        print(f"❌ Lỗi xử lý {file.name}: {e}")

print("✅ Đã watermark:", processed_files)
